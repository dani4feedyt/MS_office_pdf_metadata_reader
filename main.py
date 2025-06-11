import os, sys
from pypdf import PdfReader
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
from tkinter import filedialog
from tkinterdnd2 import DND_FILES, TkinterDnD
import tkinter as tk
import re
import json

file_path = None
result_f = ''

def resource_path(relative_path):
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

def get_pdf_metadata(f_path):
    reader = PdfReader(f_path)
    content = reader.metadata
    metadata = {
        'FileType': 'PDF document',
        'Created': content.creation_date,
        'Author': content.creator,
        'LastModified': content.modification_date,
        'ModifiedIn': content.producer_raw,
        'Title': content.title,
        'Category': content.subject,
        'Tags': content.keywords,

    }
    return metadata


def get_excel_metadata(f_path):
    workbook = load_workbook(f_path, keep_vba=False)
    metadata = {
        'FileType': 'Excel spreadsheet',
        'SheetCount': len(workbook.sheetnames),
        'SheetNames': workbook.sheetnames,
        'Created': workbook.properties.created,
        'Author': workbook.properties.creator,
        'LastModified': workbook.properties.modified,
        'LastModifiedBy': workbook.properties.last_modified_by,
        'Version': workbook.properties.version,
        'Title': workbook.properties.title,
        'Category': workbook.properties.category,
        'Tags': workbook.properties.keywords

    }
    return metadata


def get_word_metadata(f_path):
    doc = Document(f_path)
    props = doc.core_properties
    metadata = {
        'FileType': 'Word document',
        'Created': props.created,
        'Author': props.author,
        'LastModified': props.modified,
        'LastModifiedBy': props.last_modified_by,
        'Version': props.version,
        'Title': props.title,
        'Category': props.category,
        'Tags': props.keywords
    }
    return metadata


def get_ppoint_metadata(f_path):
    pres = Presentation(f_path)
    slide_names = []
    for slide in pres.slides:
        try:
            slide_title = re.sub(r'\s+', ' ', slide.shapes.title.text)
        except AttributeError:
            slide_title = '__'
        slide_names.append(slide_title)

    metadata = {
        'FileType': 'PowerPoint presentation',
        'SlideCount': len(pres.slides),
        'SlideNames': slide_names,
        'Created': pres.core_properties.created,
        'Author': pres.core_properties.author,
        'LastModified': pres.core_properties.modified,
        'LastModifiedBy': pres.core_properties.last_modified_by,
        'Version': pres.core_properties.version,
        'Title': pres.core_properties.title,
        'Category': pres.core_properties.category,
        'Tags': pres.core_properties.keywords
    }
    return metadata


def get_file_metadata(f_path):
    try:
        ext = os.path.splitext(f_path)[1].lower()
        if ext == '.pdf':
            metadata = get_pdf_metadata(f_path)
        elif ext in ['.xlsx', '.xls']:
            metadata = get_excel_metadata(f_path)
        elif ext in ['.docx', '.doc']:
            metadata = get_word_metadata(f_path)
        elif ext in ['.pptx', '.ppt']:
            metadata = get_ppoint_metadata(f_path)
        else:
            return f'Unsupported file type: {ext}'

        metadata['FileSize'] = os.stat(f_path).st_size
        metadata['FileName'] = os.path.basename(f_path)
        metadata['FilePath'] = f_path
        metadata['FileExtension'] = ext

        return metadata

    except Exception as e:
        return f'Error: {str(e)}'


def change_pic_down(event):
    button = event.widget
    if str(button.winfo_name()) == '!button':
        activated_img = tk.PhotoImage(file=resource_path('button_choose_A.png'))
        button.configure(image=activated_img)
        button.image = activated_img
        open_file()

    elif str(button.winfo_name()) == '!button2':
        activated_img = tk.PhotoImage(file=resource_path('button_extract_A.png'))
        button.configure(image=activated_img)
        button.image = activated_img
        extract_metadata(file_path)


def change_pic_up(event):
    button = event.widget
    if str(button.winfo_name()) == '!button':
        button = event.widget
        deactivated_img = tk.PhotoImage(file=resource_path('button_choose.png'))
        button.configure(image=deactivated_img)
        button.image = deactivated_img
    elif str(button.winfo_name()) == '!button2':
        activated_img = tk.PhotoImage(file=resource_path('button_extract.png'))
        button.configure(image=activated_img)
        button.image = activated_img


def text_box_write(message, textbox):
    textbox.config(state='normal')
    textbox.delete(0.0, tk.END)
    textbox.insert(tk.END, message)
    textbox.config(state='disabled')


def open_file(called=False):
    global file_path
    f_path = ''
    text_box.config(state='normal')

    if called:
        if not text_box.get(0.0, tk.END).strip("\n").endswith(('.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx')):
            text_box_write('Error, unsupported file type. Supported file types: .pdf,.xlsx,.docx', text_box)
            file_path = ''
            return

        if text_box.get(0.0, tk.END).startswith(('C:', 'D:')):
            f_path = text_box.get(0.0, tk.END).strip('\n')
    else:
        text_box.delete(0.0, tk.END)
        f_path = filedialog.askopenfilename(filetypes=[
            ('Supported Files', '*.xlsx *.xls *.docx *.doc *.pdf *.pptx *.ppt'),
            ('Excel Files', '*.xlsx *.xls'),
            ('PowerPoint Files', '*.ppt *.pptx'),
            ('Word Files', '*.docx *.doc'),
            ('PDF Files', '*.pdf')
        ])

    file_path = f_path
    text_box.delete(0.0, tk.END)
    if file_path:
        text_box_write(f'File <{os.path.basename(file_path)}> loaded successfully!', text_box)
    else:
        text_box_write('Error, please select a file!', text_box)


def handle_drop(event):
    text_box_write(re.sub(r'[{}]', '', f'{event.data}\n'), text_box)
    open_file(True)
    return event.action


def extract_metadata(f_path):
    global result_f
    text_box.config(state='normal')
    if f_path:
        result = 'Error reading file metadata'
        metadata = get_file_metadata(f_path)
        f_ext = metadata.get('FileExtension')
        if isinstance(metadata, dict):
            result = f"File Type: {metadata.get('FileType')}\n"
            result += f"File Name: {metadata.get('FileName')}\n"
            result += f"File Extension: {f_ext}\n"
            result += f"File Size: {metadata.get('FileSize')} bytes\n"
            result += f"File Path: {metadata.get('FilePath')}\n"
            if f_ext in ['.xlsx', '.xls']:
                result += f"Sheet Count: {metadata.get('SheetCount')}\n"
                result += f"Sheet Names: {metadata.get('SheetNames')}\n"
            if f_ext in ['.pptx', '.ppt']:
                result += f"Slide Count: {metadata.get('SlideCount')}\n"
                result += f"Slide Names: {metadata.get('SlideNames')}\n"
            result += f"Created: {metadata.get('Created')}\n"
            result += f"Author: {metadata.get('Author')}\n"
            result += f"Last Modified: {metadata.get('LastModified')}\n"
            if not f_ext == '.pdf':
                result += f"Last Modified By: {metadata.get('LastModifiedBy')}\n"
                result += f"Version: {'1.0' if not metadata.get('Version') else metadata.get('Version')}\n"
            else:
                result += f"Modified In: {metadata.get('ModifiedIn')}\n"
            result += f"Title: {metadata.get('Title')}\n"
            result += f"Category: {metadata.get('Category')}\n"
            result += f"Tags: {metadata.get('Tags')}\n"

        result_f = ''
        for line in result.splitlines():
            if line.endswith(' '):
                line += 'None'
            result_f += f'{line}\n'
        text_box_write(result_f, text_box)
        save_to_json(f_path)

    else:
        text_box_write('Error, invalid file or corrupted filepath!', text_box)


def save_to_json(f_path):
    filename = os.path.splitext(os.path.basename(f_path))[0]
    result_arr = result_f.split('\n')
    result_dict = {}
    for line in result_arr:
        datarow = line.split(": ")
        if datarow[0]:
            result_dict.update({datarow[0]: datarow[1]})

    json_object = json.dumps(result_dict, indent=4)
    file = filedialog.asksaveasfilename(defaultextension=".json",
                                        filetypes=[("JSON file", '*.json')],
                                        initialfile=f"{filename}.json")

    if file:
        with open(file, "w") as f:
            f.write(json_object)


root = TkinterDnD.Tk()

root.title('Dublin Core Metadata Extraction Tool')
root.geometry('700x600')
root.resizable(False, False)
true_bg = '#222831'
root.config(bg=true_bg)
root.iconbitmap(resource_path("icon.ico"))

b_img = tk.PhotoImage(file=resource_path('button_choose.png'))
open_button = tk.Button(root, image=b_img, bg=true_bg, borderwidth=0, relief='solid', activebackground=true_bg)
open_button.image = b_img
open_button.bind('<ButtonPress-1>', change_pic_down)
open_button.bind('<ButtonRelease-1>', change_pic_up)

b2_img = tk.PhotoImage(file=resource_path('button_extract.png'))
extract_button = tk.Button(root, image=b2_img, bg=true_bg, borderwidth=0, relief='solid', activebackground=true_bg)
extract_button.image = b2_img
extract_button.bind('<ButtonPress-1>', change_pic_down)
extract_button.bind('<ButtonRelease-1>', change_pic_up)

canvas = tk.Canvas(root, width=600, height=200, borderwidth=0, highlightthickness=0)
text_box = tk.Text(canvas, height=20, width=70, bg='#393E46', fg='white', font=('Consolas', 12), wrap=tk.WORD, state='disabled', borderwidth=0)
canvas.create_window(0, 0, window=text_box)
canvas.bind('<Button-1>', root.focus())
canvas.drop_target_register(DND_FILES)
canvas.dnd_bind('<<Drop>>', handle_drop)

open_button.pack(pady=20)
canvas.pack()
text_box.pack()
extract_button.pack(pady=20)

root.mainloop()
